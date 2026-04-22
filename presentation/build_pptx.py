"""
Shop Agent 课程汇报 PPTX 构建脚本
================================

用法：
    pip install python-pptx
    python build_pptx.py

功能：
    - 生成 16:9 版式、16 张幻灯片的 .pptx（含 ER 图页）
    - 与同目录 index.html 内容逐页对齐
    - 关键设计图/截图使用报告模板中的真实图片（assets/*.png）
    - 学术清爽风格（白底 + 深蓝 #1F3A68 + 蓝绿 #2E86C1 + 中灰 #566573）
    - 每页包含讲者备注

输出：
    presentation/Shop-Agent汇报.pptx
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ==================== 颜色常量 ====================
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)   # 深蓝（标题 / 主色）
COLOR_ACCENT = RGBColor(0x2E, 0x86, 0xC1)    # 蓝绿（强调）
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)      # 正文
COLOR_MUTED = RGBColor(0x56, 0x65, 0x73)     # 辅助灰
COLOR_BG_SOFT = RGBColor(0xF8, 0xFA, 0xFC)   # 浅灰底
COLOR_LINE = RGBColor(0xE5, 0xE8, 0xE8)      # 分割线
COLOR_OK = RGBColor(0x27, 0xAE, 0x60)        # 绿
COLOR_WARN = RGBColor(0xE6, 0x7E, 0x22)      # 橙
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ==================== 字号 ====================
FS_TITLE = 36
FS_SUBTITLE = 22
FS_BODY = 18
FS_SMALL = 14
FS_TINY = 11

FONT_CN = "Microsoft YaHei"

# ==================== 版式尺寸（16:9，英寸）====================
SLIDE_W = 13.333
SLIDE_H = 7.5


def _set_run_font(run, size=None, bold=False, color=None, font_name=FONT_CN):
    """统一设置 run 的字体属性，兼容中文字体。"""
    if size is not None:
        run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = font_name
    # 同时设置东亚字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)


def add_text_box(slide, left, top, width, height, text, *,
                 size=FS_BODY, bold=False, color=COLOR_TEXT,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """添加单行/多行文本框。"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        _set_run_font(run, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size=FS_BODY, color=COLOR_TEXT, line_spacing=1.2,
                bullet_color=COLOR_ACCENT):
    """添加项目符号列表。bullets 为 str 列表；若首字符是 '·'/'•' 则视为已带符号。"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing

        # bullet 点
        run_b = p.add_run()
        run_b.text = "● "
        _set_run_font(run_b, size=size, bold=True, color=bullet_color)

        # 内容
        run_t = p.add_run()
        run_t.text = item
        _set_run_font(run_t, size=size, color=color)
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=None, line=None, line_width=1.0,
             shadow=False, rounded=False):
    """添加矩形（可选圆角）。"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    if not shadow:
        shp.shadow.inherit = False
    return shp


def add_header(slide, chapter_num, chapter_name, title, subtitle=""):
    """统一页面头部：章节条 + 标题 + 分割线。"""
    # 章节条（浅底，左边蓝绿竖条）
    add_rect(slide, 0.55, 0.35, 4.5, 0.36,
             fill=COLOR_BG_SOFT, line=None, rounded=False)
    add_rect(slide, 0.55, 0.35, 0.07, 0.36, fill=COLOR_ACCENT, line=None)

    # 章节编号 + 名称
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.36),
                                  Inches(4.3), Inches(0.34))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"{chapter_num}  "
    _set_run_font(r1, size=FS_SMALL, bold=True, color=COLOR_ACCENT)
    r2 = p.add_run(); r2.text = f"/ {chapter_name}"
    _set_run_font(r2, size=FS_SMALL, bold=True, color=COLOR_PRIMARY)

    # 主标题
    add_text_box(slide, 0.55, 0.85, 12.2, 0.7, title,
                 size=FS_TITLE, bold=True, color=COLOR_PRIMARY)
    # 副标题
    if subtitle:
        add_text_box(slide, 0.55, 1.45, 12.2, 0.4, subtitle,
                     size=FS_SMALL, color=COLOR_MUTED)

    # 分割线
    add_rect(slide, 0.55, 1.95, 12.2, 0.03,
             fill=COLOR_LINE, line=None)


def add_footer(slide, page_no, total=16):
    """页脚：品牌 · 课程汇报 | 页码。"""
    add_text_box(slide, 0.55, 7.05, 8, 0.3,
                 "Shop Agent · 电商管理客服系统 · 课程汇报",
                 size=FS_TINY, color=COLOR_MUTED)
    add_text_box(slide, 11.2, 7.05, 1.6, 0.3,
                 f"{page_no:02d} / {total:02d}",
                 size=FS_TINY, color=COLOR_ACCENT, bold=True,
                 align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    """添加讲者备注。"""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text
    for p in notes_tf.paragraphs:
        for run in p.runs:
            _set_run_font(run, size=FS_SMALL, color=COLOR_TEXT)


ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")


def add_image_fit(slide, left, top, max_w, max_h, filename, *, caption=None):
    """按最大宽高等比插入图片；可选居中 caption。"""
    path = os.path.join(ASSETS_DIR, filename)
    if not os.path.exists(path):
        # 图片缺失时用占位框避免整页崩
        add_rect(slide, left, top, max_w, max_h,
                 fill=COLOR_BG_SOFT, line=COLOR_LINE, rounded=True)
        add_text_box(slide, left, top + max_h / 2 - 0.2, max_w, 0.4,
                     f"[缺失图片: {filename}]",
                     size=FS_SMALL, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
        return None

    from PIL import Image as _PILImage  # local import, 可选依赖
    try:
        with _PILImage.open(path) as im:
            iw, ih = im.size
    except Exception:
        # 无 PIL 时按 max_w 插入，由 python-pptx 等比
        pic = slide.shapes.add_picture(path, Inches(left), Inches(top),
                                       width=Inches(max_w))
        return pic

    ratio_img = iw / ih
    ratio_box = max_w / max_h
    if ratio_img >= ratio_box:
        # 宽度主导
        w = max_w
        h = max_w / ratio_img
    else:
        h = max_h
        w = max_h * ratio_img
    # 居中
    cx = left + (max_w - w) / 2
    cy = top + (max_h - h) / 2
    pic = slide.shapes.add_picture(path, Inches(cx), Inches(cy),
                                   width=Inches(w), height=Inches(h))
    if caption:
        add_text_box(slide, left, top + max_h + 0.05, max_w, 0.3,
                     caption, size=FS_TINY, color=COLOR_MUTED,
                     align=PP_ALIGN.CENTER)
    return pic


# ==================== 页面构建函数 ====================

def slide_01_cover(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景渐变（用两个矩形模拟）
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_WHITE, line=None)
    add_rect(slide, 0, 5.5, SLIDE_W, 2.0, fill=COLOR_BG_SOFT, line=None)

    # 徽标圆形
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(5.9), Inches(1.8),
                                   Inches(1.5), Inches(1.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = COLOR_PRIMARY
    badge.line.fill.background(); badge.shadow.inherit = False
    tf = badge.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "SA"
    _set_run_font(r, size=44, bold=True, color=COLOR_WHITE)

    # 主标题
    add_text_box(slide, 1, 3.6, 11.3, 0.9, "Shop Agent",
                 size=56, bold=True, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 11.3, 0.5, "电商管理客服系统",
                 size=28, bold=False, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.0, 11.3, 0.4,
                 "基于 RAG + Agent + Tool Calling 的课程大作业",
                 size=16, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # 底部元信息三列
    meta_items = [("课程", "人工智能应用"),
                  ("作者", "姓名 / 学号"),
                  ("日期", "2026 年春")]
    col_w = 3.4
    start_left = (SLIDE_W - col_w * 3) / 2
    for i, (k, v) in enumerate(meta_items):
        left = start_left + i * col_w
        add_text_box(slide, left, 6.2, col_w, 0.3, k,
                     size=13, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
        add_text_box(slide, left, 6.5, col_w, 0.36, v,
                     size=16, bold=True, color=COLOR_PRIMARY,
                     align=PP_ALIGN.CENTER)

    set_notes(slide,
              "大家好，这次课程汇报我将介绍我的大作业——Shop Agent 电商管理客服系统。"
              "它是一个基于 RAG、Agent 和 Tool Calling 的极简版智能客服，"
              "整个项目已完成后端 FastAPI 服务、Streamlit 前端，以及三阶段的 Agent Memory 压缩机制。"
              "汇报约 10 分钟，分为背景、技术、演示、总结四部分。")


def slide_02_toc(prs):
    """目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 0.55, 0.55, 12, 0.7, "汇报目录",
                 size=FS_TITLE, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 0.55, 1.2, 12, 0.35,
                 "Table of Contents",
                 size=FS_SMALL, color=COLOR_MUTED)
    add_rect(slide, 0.55, 1.7, 12.2, 0.03, fill=COLOR_LINE, line=None)

    items = [
        ("01", "项目背景与定位", "🛒 痛点分析 · 目标设定"),
        ("02", "核心技术实现", "⚙️ 架构 · RAG · Tool Calling · Memory"),
        ("03", "功能演示", "💬 登录 · 对话 · 意图 · 多轮"),
        ("04", "总结与展望", "📊 成果 · 后续工作 · Q&A"),
    ]
    # 2x2 卡片
    card_w, card_h = 5.9, 1.9
    gap_x, gap_y = 0.35, 0.35
    start_left = 0.55
    start_top = 2.15

    for i, (no, title, desc) in enumerate(items):
        row, col = divmod(i, 2)
        left = start_left + col * (card_w + gap_x)
        top = start_top + row * (card_h + gap_y)

        # 卡片背景
        add_rect(slide, left, top, card_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, line_width=1.0, rounded=True)
        # 编号
        add_text_box(slide, left + 0.3, top + 0.3, 1.4, 1.0, no,
                     size=42, bold=True, color=COLOR_ACCENT)
        # 标题
        add_text_box(slide, left + 1.7, top + 0.4, card_w - 2.0, 0.6,
                     title, size=24, bold=True, color=COLOR_PRIMARY)
        # 描述
        add_text_box(slide, left + 1.7, top + 1.0, card_w - 2.0, 0.5,
                     desc, size=14, color=COLOR_MUTED)

    add_footer(slide, 2)
    set_notes(slide,
              "汇报分为四个部分：背景与定位说明为什么要做；"
              "技术实现介绍系统的四个核心模块；"
              "功能演示展示用户真实使用流程；最后是总结与后续工作。")


def slide_03_background(prs):
    """项目背景与痛点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "01", "项目背景与定位",
               "项目背景与痛点", "Why Shop Agent")

    # 两张卡片
    # 左卡
    add_rect(slide, 0.55, 2.25, 6.05, 2.6,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 0.55, 2.25, 0.07, 2.6, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, 0.8, 2.4, 5.7, 0.4, "电商客服现状",
                 size=20, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.8, 2.85, 5.7, 2.0, [
        "SKU 规模增长，人工客服成本高、响应慢",
        "24×7 在线服务难以全人工覆盖",
        "标准化问题（订单查询 / 商品推荐）占比高",
    ], size=15)

    # 右卡
    add_rect(slide, 6.75, 2.25, 6.05, 2.6,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 6.75, 2.25, 0.07, 2.6, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, 7.0, 2.4, 5.7, 0.4, "LLM 直接应用的挑战",
                 size=20, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 7.0, 2.85, 5.7, 2.0, [
        "没有工具调用，无法访问订单数据库",
        "缺少检索增强，无法给出真实商品推荐",
        "上下文暴涨，Token 与延迟线性恶化",
    ], size=15)

    # 解决方案条
    add_rect(slide, 0.55, 5.15, 12.25, 1.5,
             fill=COLOR_BG_SOFT, line=COLOR_PRIMARY, line_width=1.2, rounded=True)
    add_rect(slide, 0.55, 5.15, 0.1, 1.5, fill=COLOR_PRIMARY, line=None)
    add_text_box(slide, 0.85, 5.3, 12.0, 0.4, "→ 解决方案",
                 size=20, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 0.85, 5.8, 11.9, 0.75,
                 "在大模型之上叠加三层能力：Tool Calling 接数据库、"
                 "RAG 接知识库、Agent Memory 压长上下文。",
                 size=16, color=COLOR_TEXT)

    add_footer(slide, 3)
    set_notes(slide,
              "电商客服有三个典型痛点：成本高、7×24 覆盖难、标准化问题多。"
              "纯 LLM 无法访问真实订单数据，也无法给出准确商品推荐，"
              "而且对话一长 Token 就爆了。所以我们用 Tool Calling + RAG + Memory 三件套来解决。")


def slide_04_positioning(prs):
    """项目定位与目标"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "01", "项目背景与定位",
               "项目定位与目标", "Goal & Scope")

    # 左卡 · 定位
    add_rect(slide, 0.55, 2.25, 6.05, 4.3,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 0.55, 2.25, 0.07, 4.3, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, 0.8, 2.4, 5.7, 0.4, "项目定位",
                 size=20, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 0.8, 2.9, 5.7, 1.2,
                 "面向中小型电商的极简可落地智能客服雏形，"
                 "覆盖从 API 到前端的完整链路。",
                 size=15, color=COLOR_TEXT)
    # 标签
    tags = ["大模型对话", "工具调用", "检索增强", "Agent 记忆"]
    tag_left = 0.8
    tag_top = 4.2
    for i, tag in enumerate(tags):
        row, col = divmod(i, 2)
        lx = tag_left + col * 2.6
        ly = tag_top + row * 0.55
        add_rect(slide, lx, ly, 2.4, 0.42,
                 fill=COLOR_BG_SOFT, line=COLOR_LINE, rounded=True)
        add_text_box(slide, lx, ly + 0.05, 2.4, 0.32, tag,
                     size=13, bold=True, color=COLOR_ACCENT,
                     align=PP_ALIGN.CENTER)

    # 右卡 · 目标
    add_rect(slide, 6.75, 2.25, 6.05, 4.3,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 6.75, 2.25, 0.07, 4.3, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, 7.0, 2.4, 5.7, 0.4, "具体目标（全部完成）",
                 size=20, bold=True, color=COLOR_PRIMARY)
    goals = [
        "✅ 基础对话 API（/api/chat）",
        "✅ 工具调用：订单 & 物流查询",
        "✅ RAG：基于产品向量库的推荐",
        "✅ 意图路由：自动选择合适 Agent",
        "✅ Agent Memory：对话历史压缩",
        "✅ Streamlit 前端 + JWT 认证",
    ]
    add_bullets(slide, 7.0, 2.95, 5.7, 3.5, goals,
                size=15, bullet_color=COLOR_OK)

    add_footer(slide, 4)
    set_notes(slide,
              "定位是一个小而全的可落地原型，重点不在规模，"
              "而在串通 RAG + Agent + Memory 的完整链路。"
              "右边是 6 个主要目标，全部已经实现，后面会逐项展开。")


def slide_05_tech_stack(prs):
    """技术栈总览"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "技术栈总览", "Tech Stack")

    stack = [
        ("FastAPI", "异步 Web 框架"),
        ("DeepSeek API", "大模型（OpenAI 兼容）"),
        ("SiliconFlow", "Embedding 服务"),
        ("ChromaDB", "向量数据库"),
        ("SQLAlchemy", "ORM / SQLite"),
        ("Pydantic v2", "数据校验 & 配置"),
        ("Streamlit", "前端聊天界面"),
        ("Uvicorn", "ASGI 服务器"),
    ]

    # 4x2 网格
    card_w, card_h = 2.9, 1.15
    gap_x, gap_y = 0.18, 0.2
    start_left = 0.55
    start_top = 2.3
    for i, (name, role) in enumerate(stack):
        row, col = divmod(i, 4)
        left = start_left + col * (card_w + gap_x)
        top = start_top + row * (card_h + gap_y)
        add_rect(slide, left, top, card_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        add_text_box(slide, left, top + 0.22, card_w, 0.45, name,
                     size=17, bold=True, color=COLOR_PRIMARY,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, left, top + 0.7, card_w, 0.4, role,
                     size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # 工程化支撑条
    add_rect(slide, 0.55, 5.35, 12.25, 1.3,
             fill=COLOR_BG_SOFT, line=COLOR_LINE, rounded=True)
    add_text_box(slide, 0.85, 5.5, 11.8, 0.4,
                 "工程化支撑", size=18, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 0.85, 6.0, 11.8, 0.55,
                 "启动配置校验  ·  请求日志中间件  ·  速率限制中间件  ·  "
                 "全局异常处理  ·  psutil 监控端点",
                 size=14, color=COLOR_TEXT)

    add_footer(slide, 5)
    set_notes(slide,
              "技术栈全部是 Python 生态：FastAPI 做后端，DeepSeek 做大模型，"
              "SiliconFlow 做向量化，ChromaDB 做向量检索，SQLite 做订单库，"
              "Streamlit 做前端。下面一排是工程化组件：启动校验、日志、限流、异常、监控。")


def slide_06_architecture(prs):
    """整体架构图 · 使用报告图 7（真实架构图）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "系统整体架构", "Layered Architecture · 四层解耦")

    # 左：真实架构图
    add_image_fit(slide, 0.55, 2.15, 5.6, 4.8,
                  "diagram-architecture.png")

    # 右：四层说明卡
    cards = [
        ("① 展示层 Presentation",
         "微信小程序（正式演示） + Streamlit（开发联调）"),
        ("② 接口与网关层 API Gateway",
         "FastAPI REST + 鉴权 / 限流 / 日志中间件"),
        ("③ 智能业务层 AI & Agent",
         "RouterAgent → Order · RAG · General，旁挂 AgentMemory"),
        ("④ 数据与持久层 Data",
         "SQLite 关系型业务数据 + ChromaDB 商品向量库"),
    ]
    card_left = 6.4
    card_w = 6.4
    card_h = 1.08
    gap = 0.12
    start_top = 2.25
    for i, (t, d) in enumerate(cards):
        ly = start_top + i * (card_h + gap)
        add_rect(slide, card_left, ly, card_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        add_rect(slide, card_left, ly, 0.07, card_h, fill=COLOR_ACCENT, line=None)
        add_text_box(slide, card_left + 0.22, ly + 0.12, card_w - 0.3, 0.42, t,
                     size=15, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, card_left + 0.22, ly + 0.58, card_w - 0.3, 0.46, d,
                     size=12, color=COLOR_TEXT)

    add_footer(slide, 6)
    set_notes(slide,
              "系统采用\"分层后端 + 多 Agent 协同 + 双前端验证\"的方案，共四层："
              "展示层有微信小程序和 Streamlit 两套前端；API 网关层统一通过 FastAPI 并加鉴权、限流、日志中间件；"
              "智能层由 RouterAgent 分发到 Order、RAG、General 三类 Agent，旁挂 AgentMemory 做记忆压缩；"
              "数据层是 SQLite 关系库加 ChromaDB 向量库。这张图在我们的报告中是\"图 3·系统架构图\"。")


def slide_07_tool_calling(prs):
    """Tool Calling · 使用报告图 2（UML 时序图）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "核心技术（一）Tool Calling", "订单查询 Agent · UML 时序")

    # 左：真实时序图
    add_image_fit(slide, 0.55, 2.2, 7.3, 4.65,
                  "diagram-tool-calling.png")

    # 右上：7 步流程
    code_left = 8.1
    add_rect(slide, code_left, 2.2, 4.75, 2.2,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, code_left, 2.2, 0.07, 2.2, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, code_left + 0.2, 2.3, 4.5, 0.4,
                 "7 步调用流程", size=15, bold=True, color=COLOR_PRIMARY)
    steps = [
        "① 用户发起自然语言请求",
        "② OrderAgent 判意图、匹配 Tools",
        "③ 触发 Tool Calling，参数校验",
        "④ OrderService 执行 SQL 查询",
        "⑤ SQLite 返回业务数据",
        "⑥ 结果序列化为 ToolMessage 回喂",
        "⑦ LLM 结合上下文生成最终答复",
    ]
    tb = slide.shapes.add_textbox(Inches(code_left + 0.2), Inches(2.75),
                                  Inches(4.5), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    for i, s in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        r = p.add_run(); r.text = s
        _set_run_font(r, size=11, color=COLOR_TEXT)

    # 右下：工具定义代码
    add_rect(slide, code_left, 4.55, 4.75, 2.35,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, code_left, 4.55, 0.07, 2.35, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, code_left + 0.2, 4.65, 4.5, 0.4,
                 "工具定义（OpenAI 兼容）", size=15, bold=True, color=COLOR_PRIMARY)
    code_text = ('tools = [{\n'
                 '  "type": "function",\n'
                 '  "function": {\n'
                 '    "name": "get_order_status",\n'
                 '    "description": "查询订单物流状态",\n'
                 '    "parameters": { ...order_number... }\n'
                 '  }\n'
                 '}]')
    add_text_box(slide, code_left + 0.2, 5.1, 4.5, 1.75, code_text,
                 size=10, color=COLOR_TEXT)

    add_footer(slide, 7)
    set_notes(slide,
              "Tool Calling 用的是 OpenAI 函数调用协议。左侧是系统的 UML 时序图，"
              "完整展示了用户 → OrderAgent（LLM） → OrderService（工具调度层） → SQLite 的 7 步交互："
              "模型先做语义理解和意图判断，触发 tool_calls 后由后端执行真实 SQL 查询，"
              "再把结果以 ToolMessage 角色回喂给模型，最终生成自然语言回复。")


def slide_08_rag(prs):
    """RAG 产品推荐 · 使用报告图 4（两阶段流程）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "核心技术（二）RAG 产品推荐", "检索增强生成 · 两阶段")

    # 左：真实流程图（竖图）
    add_image_fit(slide, 0.55, 2.15, 4.8, 4.85,
                  "diagram-rag-flow.png")

    # 右：三张说明卡
    descs = [
        ("一、向量检索阶段",
         "用户 Query → Embedding 模型 → Query 向量 → ChromaDB 计算相似度 → Top-K=5 候选商品。"),
        ("二、增强生成阶段",
         "Top-K 结果注入 Prompt，与原始意图拼接后交由 DeepSeek LLM，生成个性化推荐话术。"),
        ("数据与依赖",
         "data/products.csv 30 款商品离线入库；Embedding 由 SiliconFlow 提供（OpenAI 兼容）。"),
    ]
    right_left = 5.65
    right_w = 7.2
    card_h = 1.5
    gap = 0.18
    start_top = 2.3
    for i, (t, d) in enumerate(descs):
        ly = start_top + i * (card_h + gap)
        add_rect(slide, right_left, ly, right_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        color = COLOR_PRIMARY if i == len(descs) - 1 else COLOR_ACCENT
        add_rect(slide, right_left, ly, 0.08, card_h, fill=color, line=None)
        add_text_box(slide, right_left + 0.25, ly + 0.18, right_w - 0.35, 0.5, t,
                     size=18, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, right_left + 0.25, ly + 0.75, right_w - 0.35, 0.7, d,
                     size=13, color=COLOR_TEXT)

    add_footer(slide, 8)
    set_notes(slide,
              "RAG 在我们系统中分两个阶段：左图上半部分是向量检索，把用户 query 转成向量后"
              "在 ChromaDB 里做相似度检索，召回 Top-K 商品；下半部分是增强生成，"
              "把召回结果和原始意图一起拼进 Prompt 交给大模型生成推荐话术。"
              "相比关键词匹配，向量检索能处理语义近义表达，鲁棒性更高。")


def slide_09_router(prs):
    """Router Agent"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "核心技术（三）Router Agent", "意图路由")

    # 左 · 三分类
    add_rect(slide, 0.55, 2.25, 5.7, 4.3,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 0.55, 2.25, 0.07, 4.3, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, 0.8, 2.4, 5.3, 0.4,
                 "三分类意图识别", size=20, bold=True, color=COLOR_PRIMARY)

    categories = [
        ("order_query", "订单 / 物流 / 快递相关"),
        ("product_rec", "商品推荐 / 选购咨询"),
        ("general", "寒暄 / 政策 / 其它"),
    ]
    for i, (k, v) in enumerate(categories):
        ky = 2.95 + i * 0.85
        # 标签
        add_rect(slide, 0.9, ky, 2.3, 0.55, fill=COLOR_PRIMARY, line=None, rounded=True)
        add_text_box(slide, 0.9, ky + 0.1, 2.3, 0.35, k,
                     size=13, bold=True, color=COLOR_WHITE,
                     align=PP_ALIGN.CENTER)
        # 值
        add_rect(slide, 3.3, ky, 2.8, 0.55,
                 fill=COLOR_BG_SOFT, line=COLOR_LINE, rounded=True)
        add_text_box(slide, 3.3, ky + 0.1, 2.8, 0.35, v,
                     size=12, color=COLOR_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, 0.8, 5.75, 5.3, 0.6,
                 "由 router_agent.py 统一调度，决定走 Order / RAG / 通用对话。",
                 size=13, color=COLOR_MUTED)

    # 右上 · 调用时序
    add_rect(slide, 6.45, 2.25, 6.35, 2.1,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 6.45, 2.25, 0.07, 2.1, fill=COLOR_ACCENT, line=None)
    add_text_box(slide, 6.7, 2.4, 6.0, 0.4, "调用时序",
                 size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 6.7, 2.85, 6.0, 1.5, [
        "用户消息 + 历史 → Router Agent",
        "LLM 输出意图标签",
        "分发到对应子 Agent 执行",
        "统一回包（含 intent 字段）",
    ], size=13, line_spacing=1.15)

    # 右下 · 优化点
    add_rect(slide, 6.45, 4.5, 6.35, 2.05,
             fill=COLOR_WHITE, line=COLOR_WARN, rounded=True)
    add_rect(slide, 6.45, 4.5, 0.07, 2.05, fill=COLOR_WARN, line=None)
    add_text_box(slide, 6.7, 4.6, 6.0, 0.4, "已识别优化点",
                 size=18, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 6.7, 5.05, 6.0, 1.4,
                 "额外的意图分类 LLM 调用会增加延迟与成本，"
                 "已在 to-fix.md 标记，后续可用规则预过滤 + 单次 tool_choice 合并。",
                 size=13, color=COLOR_TEXT)

    add_footer(slide, 9)
    set_notes(slide,
              "Router Agent 做的是三分类：订单查询、产品推荐、通用对话。"
              "它用 LLM 识别意图后分发给对应子 Agent，前端也会显示识别出的意图标签。"
              "已识别到的优化点是分类多了一次 LLM 调用，后续可以合并或用规则预过滤。")


def slide_10_memory(prs):
    """Agent Memory 三阶段"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "核心技术（四）Agent Memory 三阶段", "上下文压缩")

    # 三卡片
    cards = [
        ("阶段 1 · Buffer 策略", "规则截断", COLOR_PRIMARY, [
            "最近 3 轮完整保留",
            "第 4-8 轮仅保留摘要",
            "第 9 轮及更早删除",
            "Tool / RAG 结果标记为关键",
        ], "对应：agent_memory.py"),
        ("阶段 2 · 场景感知", "结构化提取", COLOR_ACCENT, [
            "自动识别对话模式",
            "抽取订单号 / 状态 / 物流",
            "抽取商品 / 类目 / 预算",
            "Pattern detection ≥ 95%",
        ], "结构化保存关键字段"),
        ("阶段 3 · LLM 压缩", "智能摘要", COLOR_PRIMARY, [
            "超阈值触发 LLM 压缩",
            "早期多轮 → 50 字摘要",
            "摘要 + 最近 N 轮拼装",
            "压缩耗时 P99 < 1s",
        ], "对应：summary_service.py"),
    ]
    total_w = 12.25
    card_w = (total_w - 0.4) / 3
    card_top = 2.3
    card_h = 3.2
    for i, (title, sub, color, bullets, footnote) in enumerate(cards):
        lx = 0.55 + i * (card_w + 0.2)
        # 外框
        add_rect(slide, lx, card_top, card_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        # 顶部色条
        add_rect(slide, lx, card_top, card_w, 0.55,
                 fill=color, line=None, rounded=True)
        # 覆盖下半圆角
        add_rect(slide, lx, card_top + 0.3, card_w, 0.25,
                 fill=color, line=None)
        add_text_box(slide, lx, card_top + 0.08, card_w, 0.4, title,
                     size=15, bold=True, color=COLOR_WHITE,
                     align=PP_ALIGN.CENTER)
        # 副标题
        add_text_box(slide, lx, card_top + 0.7, card_w, 0.4, sub,
                     size=18, bold=True, color=COLOR_PRIMARY,
                     align=PP_ALIGN.CENTER)
        # bullets
        add_bullets(slide, lx + 0.2, card_top + 1.3, card_w - 0.3, 1.5,
                    bullets, size=12, line_spacing=1.15)
        # 备注
        add_text_box(slide, lx + 0.2, card_top + 2.85, card_w - 0.3, 0.3,
                     footnote, size=11, color=COLOR_MUTED)

    # 收益表
    table_top = 5.65
    add_text_box(slide, 0.55, table_top, 12.25, 0.3,
                 "预期收益（来源：AGENT_MEMORY_PLAN.md，设计目标非实测）",
                 size=12, color=COLOR_MUTED)

    # 收益四列
    bench = [
        ("10 轮 Token", "2000 → 600", "↓ 70%"),
        ("30 轮 Token", "6000 → 1200", "↓ 80%"),
        ("平均响应时间", "2.5s → 1.5s", "↓ 40%"),
        ("API 成本/100", "$0.50 → $0.12", "↓ 76%"),
    ]
    col_w = (total_w - 0.45) / 4
    btop = 6.0
    for i, (k, v, d) in enumerate(bench):
        lx = 0.55 + i * (col_w + 0.15)
        add_rect(slide, lx, btop, col_w, 0.72,
                 fill=COLOR_BG_SOFT, line=COLOR_LINE, rounded=True)
        add_text_box(slide, lx, btop + 0.05, col_w, 0.3, k,
                     size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
        add_text_box(slide, lx, btop + 0.3, col_w, 0.2, v,
                     size=12, color=COLOR_TEXT, align=PP_ALIGN.CENTER)
        add_text_box(slide, lx, btop + 0.5, col_w, 0.25, d,
                     size=13, bold=True, color=COLOR_OK,
                     align=PP_ALIGN.CENTER)

    add_footer(slide, 10)
    set_notes(slide,
              "Memory 是本项目的亮点之一，分三阶段："
              "阶段 1 用规则 Buffer 策略保留最近 3 轮、压缩 4-8 轮；"
              "阶段 2 加入场景识别，结构化保存订单号、商品偏好等关键字段；"
              "阶段 3 用 LLM 对早期历史生成 50 字摘要。"
              "需要强调：下表数字是设计目标，来源 AGENT_MEMORY_PLAN.md，不是实测值。")


def slide_11_er(prs):
    """数据模型 ER 图 · 使用报告图 5"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "数据模型设计", "ER Diagram · 四张核心表")

    # 左：ER 图
    add_image_fit(slide, 0.55, 2.15, 6.2, 4.85,
                  "diagram-er.png")

    # 右：四张表卡
    tables = [
        ("users · 用户",
         "用户名 / 邮箱唯一索引，存储 hashed_password；一对多关联 orders。"),
        ("orders · 订单",
         "订单号、物流单号、收货地址、状态、金额、创建 / 发货 / 送达时间戳。"),
        ("products · 商品",
         "名称 / 描述 / 类别 / 价格 / 库存，支撑 RAG 向量库同步。"),
        ("order_items · 订单明细",
         "订单 × 商品桥表，保留下单时的单价与小计，保证历史一致性。"),
    ]
    right_left = 7.05
    right_w = 5.8
    card_h = 1.08
    gap = 0.12
    start_top = 2.25
    for i, (t, d) in enumerate(tables):
        ly = start_top + i * (card_h + gap)
        add_rect(slide, right_left, ly, right_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        add_rect(slide, right_left, ly, 0.07, card_h, fill=COLOR_ACCENT, line=None)
        add_text_box(slide, right_left + 0.22, ly + 0.12, right_w - 0.3, 0.42, t,
                     size=15, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, right_left + 0.22, ly + 0.58, right_w - 0.3, 0.46, d,
                     size=11, color=COLOR_TEXT)

    add_footer(slide, 11)
    set_notes(slide,
              "数据层一共四张表：users、orders、products 以及桥表 order_items。"
              "users 到 orders 是一对多（一个用户可以有多个订单），"
              "orders 与 products 通过 order_items 多对多关联。"
              "关键字段加了索引以支撑订单号、快递号、商品维度的高效检索。"
              "这张图在报告中是\"图 5·ER 图\"。")


def slide_12_highlights(prs):
    """系统亮点小结（原 slide_11_highlights）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "核心技术实现",
               "系统亮点小结", "Highlights")

    items = [
        ("① 分层清晰的架构",
         "API / Agent / Service / Database 四层解耦，模块可独立测试与替换。"),
        ("② 能力正交的 Agent",
         "Router / Order / RAG 各司其职，扩展新场景只需新增 Agent，不动路由逻辑。"),
        ("③ 三阶段记忆压缩",
         "Buffer → 场景感知 → LLM 压缩，从规则到智能，预期 Token 大幅下降。"),
        ("④ 工程化完善",
         "启动配置校验、请求日志、速率限制、JWT 认证、psutil 监控端点，贴近生产。"),
    ]
    card_w = 5.95
    card_h = 1.95
    gap_x = 0.35
    gap_y = 0.3
    for i, (t, d) in enumerate(items):
        row, col = divmod(i, 2)
        lx = 0.55 + col * (card_w + gap_x)
        ly = 2.3 + row * (card_h + gap_y)
        add_rect(slide, lx, ly, card_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        add_rect(slide, lx, ly, 0.07, card_h, fill=COLOR_ACCENT, line=None)
        add_text_box(slide, lx + 0.25, ly + 0.2, card_w - 0.3, 0.5, t,
                     size=18, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, lx + 0.25, ly + 0.85, card_w - 0.3, 1.0, d,
                     size=14, color=COLOR_TEXT)

    add_footer(slide, 12)
    set_notes(slide,
              "四大亮点：分层架构、正交 Agent、三阶段压缩、工程化配套。"
              "这些不仅是课程作业的完成度体现，也是工程上可以继续扩展的基础。")


def slide_13_demo1(prs):
    """功能演示一 · 使用真实截图（小程序登录 + 推荐结果）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "03", "功能演示",
               "功能演示（一）登录 & 商品推荐", "微信小程序 · 真机截图")

    # 左：登录截图
    add_image_fit(slide, 0.8, 2.2, 4.5, 4.4,
                  "shot-miniapp-login.png",
                  caption="登录 / 注册页（JWT 鉴权）")

    # 右：RAG 推荐结果截图
    add_image_fit(slide, 5.8, 2.2, 7.05, 4.4,
                  "shot-rag-result.png",
                  caption="RAG 推荐结果：用户问\"推荐 400 元键盘\" → ikbc C87 推荐话术")

    add_footer(slide, 13)
    set_notes(slide,
              "这一页是微信小程序端的真实运行截图。左图是登录页，携带 JWT Bearer Token "
              "调用后端统一接口；右图是商品推荐场景——用户提问\"推荐 400 元左右的键盘\"，"
              "系统经 RouterAgent 识别为 product_recommend 意图，RAGAgent 在 ChromaDB 中"
              "检索召回，再由 LLM 生成结构化的推荐话术：首选、价格、特点、适合人群与其他选项对比。")


def slide_14_demo2(prs):
    """功能演示二 · Streamlit 真实截图 + 场景要点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "03", "功能演示",
               "功能演示（二）开发端 & 场景", "Streamlit 真实截图 · 三场景")

    # 左：Streamlit 截图
    add_image_fit(slide, 0.55, 2.15, 7.2, 4.6,
                  "shot-streamlit.png",
                  caption="Streamlit 联调页：登录态 + 完整 RAG 对话 + 识别意图/使用工具 标签")

    # 右：三张场景卡 + 要点
    right_left = 8.05
    right_w = 4.8
    scenarios = [
        ("A · 订单查询（Tool Calling）",
         "ORD123 到哪了？ → Router → OrderAgent → get_order_status → SQLite。"),
        ("B · 商品推荐（RAG）",
         "300 元左右的键盘 → RAGAgent → Top-K → LLM 生成推荐。"),
        ("C · 多轮压缩（Memory）",
         "早期历史触发 LLM 摘要，Token 占用维持稳定区间。"),
    ]
    card_h = 1.1
    gap = 0.12
    start_top = 2.2
    for i, (t, d) in enumerate(scenarios):
        ly = start_top + i * (card_h + gap)
        add_rect(slide, right_left, ly, right_w, card_h,
                 fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
        add_rect(slide, right_left, ly, 0.07, card_h, fill=COLOR_ACCENT, line=None)
        add_text_box(slide, right_left + 0.22, ly + 0.12, right_w - 0.3, 0.4, t,
                     size=13, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, right_left + 0.22, ly + 0.5, right_w - 0.3, 0.58, d,
                     size=11, color=COLOR_TEXT)

    # 可视化要点卡（右下）
    ly = start_top + 3 * (card_h + gap)
    add_rect(slide, right_left, ly, right_w, 2.55,
             fill=COLOR_WHITE, line=COLOR_PRIMARY, line_width=1.2, rounded=True)
    add_rect(slide, right_left, ly, 0.09, 2.55, fill=COLOR_PRIMARY, line=None)
    add_text_box(slide, right_left + 0.22, ly + 0.12, right_w - 0.3, 0.4,
                 "可视化要点", size=14, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, right_left + 0.22, ly + 0.55, right_w - 0.3, 2.0,
                 "截图下方的「识别意图」与「使用工具」标签，"
                 "能直观证明 RouterAgent 做了意图分类、RAGAgent 真实调用了底层检索，"
                 "不是凭空生成。",
                 size=11, color=COLOR_TEXT)

    add_footer(slide, 14)
    set_notes(slide,
              "Streamlit 是开发阶段的联调前端，不作为正式演示界面。"
              "这张截图里重点强调两处标签：「识别意图：product_recommend」和「使用工具：product_search」，"
              "它们证明 RouterAgent 确实做了意图分类，RAGAgent 确实调用了底层检索能力。"
              "三大场景——订单、推荐、多轮，分别对应 Tool Calling、RAG、Memory 三项核心能力。")


def slide_15_summary(prs):
    """总结与后续（原 slide_14_summary，页码 15）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "04", "总结与展望",
               "项目总结与后续工作", "Summary & Roadmap")

    # 左 · 已完成
    add_rect(slide, 0.55, 2.3, 6.05, 4.4,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 0.55, 2.3, 0.08, 4.4, fill=COLOR_OK, line=None)
    add_text_box(slide, 0.8, 2.45, 5.5, 0.45, "✅ 已完成",
                 size=20, bold=True, color=COLOR_OK)
    done = [
        "FastAPI 完整后端 + 5 个聊天端点",
        "Tool Calling 订单查询链路",
        "ChromaDB + SiliconFlow 的 RAG 推荐",
        "Router Agent 意图路由",
        "Agent Memory 阶段 1 / 2 / 3",
        "双前端（微信小程序 + Streamlit）+ JWT 认证",
        "启动配置校验 + 日志/限流中间件",
    ]
    add_bullets(slide, 0.8, 3.0, 5.6, 3.6, done,
                size=14, bullet_color=COLOR_OK, line_spacing=1.3)

    # 右 · 后续工作
    add_rect(slide, 6.75, 2.3, 6.05, 4.4,
             fill=COLOR_WHITE, line=COLOR_LINE, rounded=True)
    add_rect(slide, 6.75, 2.3, 0.08, 4.4, fill=COLOR_WARN, line=None)
    add_text_box(slide, 7.0, 2.45, 5.5, 0.45, "⏳ 后续工作",
                 size=20, bold=True, color=COLOR_WARN)
    todo = [
        "Memory 阶段 4：会话持久化（ChatSession 表）",
        "JWT 签名验证 + 密码 bcrypt/argon2 升级",
        "意图分类合并为单次 tool_choice，降低延迟",
        "产品数据源（CSV 与 DB）统一同步机制",
        "流式响应 / WebSocket 实时聊天",
        "Dockerfile + 部署脚本 + 测试补全",
    ]
    add_bullets(slide, 7.0, 3.0, 5.6, 3.6, todo,
                size=14, bullet_color=COLOR_WARN, line_spacing=1.3)

    add_footer(slide, 15)
    set_notes(slide,
              "左边是已交付的功能清单，对应 todo.md 打勾的部分以及 AGENT_MEMORY_PLAN.md 的阶段 1-3；"
              "右边是 to-fix.md 中识别到的待完善项，优先级较高的是会话持久化和 JWT 签名升级。"
              "这既是课程作业的收尾，也是继续演进的蓝图。")


def slide_16_qa(prs):
    """Q&A 致谢（原 slide_15_qa，页码 16）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_BG_SOFT, line=None)

    add_text_box(slide, 0, 1.8, SLIDE_W, 1.8, "Q & A",
                 size=120, bold=True, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, 0, 3.8, SLIDE_W, 0.6,
                 "感谢聆听，欢迎提问与指导",
                 size=22, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0, 4.8, SLIDE_W, 0.9, "谢  谢",
                 size=54, bold=True, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, 0, 6.3, SLIDE_W, 0.3,
                 "项目仓库：Shop-Agent · main 分支",
                 size=13, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 16)
    set_notes(slide,
              "汇报到此结束，感谢老师和同学们的聆听。"
              "欢迎针对架构、RAG 召回策略、Memory 压缩效果等任意方向提问。")


# ---- 旧函数保留以便兼容（已不被 build() 使用）----

def slide_11_highlights(prs):
    """[已废弃] 使用 slide_12_highlights 替代。"""
    slide_12_highlights(prs)


def slide_12_demo1(prs):
    """[已废弃] 使用 slide_13_demo1 替代。"""
    slide_13_demo1(prs)


def slide_13_demo2(prs):
    """[已废弃] 使用 slide_14_demo2 替代。"""
    slide_14_demo2(prs)


def slide_14_summary(prs):
    """[已废弃] 使用 slide_15_summary 替代。"""
    slide_15_summary(prs)


def slide_15_qa(prs):
    """[已废弃] 使用 slide_16_qa 替代。"""
    slide_16_qa(prs)


# ==================== 主流程 ====================

def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_01_cover(prs)
    slide_02_toc(prs)
    slide_03_background(prs)
    slide_04_positioning(prs)
    slide_05_tech_stack(prs)
    slide_06_architecture(prs)
    slide_07_tool_calling(prs)
    slide_08_rag(prs)
    slide_09_router(prs)
    slide_10_memory(prs)
    slide_11_er(prs)
    slide_12_highlights(prs)
    slide_13_demo1(prs)
    slide_14_demo2(prs)
    slide_15_summary(prs)
    slide_16_qa(prs)

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "Shop-Agent汇报.pptx")
    prs.save(out_path)
    print(f"[OK] PPTX saved: {out_path}")


if __name__ == "__main__":
    build()
